"""
UNIDO AI Chatbot — Slide Deck Generator
Generates a professional PowerPoint presentation with UNIDO branding.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Palette ───────────────────────────────────────────────────────────────────
UNIDO_BLUE    = RGBColor(0x00, 0x66, 0xB3)
UNIDO_DARK    = RGBColor(0x0A, 0x1F, 0x3D)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG      = RGBColor(0xF4, 0xF6, 0xFA)
SLATE_700     = RGBColor(0x33, 0x41, 0x55)
SLATE_500     = RGBColor(0x64, 0x74, 0x8B)
ACCENT_GREEN  = RGBColor(0x05, 0x96, 0x69)
ACCENT_RED    = RGBColor(0xDC, 0x26, 0x26)
ACCENT_AMBER  = RGBColor(0xD9, 0x77, 0x06)
ACCENT_PURPLE = RGBColor(0x7C, 0x3A, 0xED)

# Slide dimensions (widescreen 16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helper functions ──────────────────────────────────────────────────────────

def add_gradient_bg(slide, color1=UNIDO_DARK, color2=UNIDO_BLUE):
    """Add a solid dark background (gradient not easily supported in python-pptx)."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color1


def add_light_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG


def add_white_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_shape(slide, left, top, width, height, fill_color, corner_radius=Emu(100000)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, font_size=14, color=SLATE_700, bold=False, alignment=PP_ALIGN.LEFT):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return p


def add_paragraph(tf, text, font_size=14, color=SLATE_700, bold=False, space_before=Pt(4), space_after=Pt(2), level=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_before = space_before
    p.space_after = space_after
    p.level = level
    return p


def add_bullet(tf, text, font_size=13, color=SLATE_700, bold_prefix="", space_before=Pt(3)):
    p = tf.add_paragraph()
    p.space_before = space_before
    p.space_after = Pt(1)
    if bold_prefix:
        run_b = p.add_run()
        run_b.text = bold_prefix + " "
        run_b.font.size = Pt(font_size)
        run_b.font.color.rgb = color
        run_b.font.bold = True
        run_n = p.add_run()
        run_n.text = text
        run_n.font.size = Pt(font_size)
        run_n.font.color.rgb = color
    else:
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
    return p


# ── Slide Header Bar ─────────────────────────────────────────────────────────

def add_header_bar(slide, slide_number, total_slides):
    """Blue strip at the top with UNIDO branding and slide counter."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.85)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = UNIDO_DARK
    bar.line.fill.background()

    # Brand text
    txb = add_text_box(slide, Inches(0.6), Inches(0.15), Inches(5), Inches(0.55))
    set_text(txb.text_frame, "UNIDO AI Chatbot", font_size=18, color=WHITE, bold=True)

    # Slide counter
    txb2 = add_text_box(slide, Inches(10.5), Inches(0.15), Inches(2.5), Inches(0.55))
    set_text(txb2.text_frame, f"{slide_number} / {total_slides}", font_size=14, color=RGBColor(0xA0, 0xC0, 0xE0), alignment=PP_ALIGN.RIGHT)


def add_footer_bar(slide):
    """Thin footer line."""
    txb = add_text_box(slide, Inches(0.6), Inches(7.0), Inches(12), Inches(0.4))
    set_text(txb.text_frame, "UNIDO Careers RAG Chatbot  •  MouriTech  •  April 2026", font_size=9, color=SLATE_500, alignment=PP_ALIGN.CENTER)


# ── Content Slide Template ────────────────────────────────────────────────────

def create_content_slide(slide_num, total, title, subtitle, bullets, accent_color=UNIDO_BLUE):
    """Standard content slide with title, subtitle, accent stripe, and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_white_bg(slide)
    add_header_bar(slide, slide_num, total)
    add_footer_bar(slide)

    # Accent stripe left
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.15), Inches(0.08), Inches(0.65)
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent_color
    stripe.line.fill.background()

    # Title
    txb_title = add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5))
    set_text(txb_title.text_frame, title, font_size=28, color=UNIDO_DARK, bold=True)

    # Subtitle
    txb_sub = add_text_box(slide, Inches(0.8), Inches(1.65), Inches(11), Inches(0.4))
    set_text(txb_sub.text_frame, subtitle, font_size=14, color=SLATE_500)

    # Bullets area
    txb_body = add_text_box(slide, Inches(0.8), Inches(2.25), Inches(11.5), Inches(4.5))
    tf = txb_body.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        # Parse bold prefix (text before " — " or ": ")
        bold_part = ""
        rest = bullet
        for sep in [" — ", ": "]:
            if sep in bullet:
                parts = bullet.split(sep, 1)
                bold_part = "•  " + parts[0] + sep.rstrip()
                rest = parts[1]
                break
        else:
            bold_part = ""
            rest = "•  " + bullet

        sp = Pt(10) if i == 0 else Pt(7)
        if bold_part:
            add_bullet(tf, rest, font_size=13, color=SLATE_700, bold_prefix=bold_part, space_before=sp)
        else:
            add_bullet(tf, rest, font_size=13, color=SLATE_700, space_before=sp)

    return slide


# ── Title Slide ───────────────────────────────────────────────────────────────

def create_title_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide)

    # Main title
    txb = add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2))
    tf = txb.text_frame
    tf.word_wrap = True
    set_text(tf, "UNIDO AI Chatbot", font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_paragraph(tf, "System Architecture & Technical Deep Dive", font_size=22, color=RGBColor(0xA0, 0xC0, 0xE0), space_before=Pt(12))
    tf.paragraphs[1].alignment = PP_ALIGN.CENTER

    # Divider line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(3.7), Inches(4.3), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = UNIDO_BLUE
    line.line.fill.background()

    # Subtitle
    txb2 = add_text_box(slide, Inches(1), Inches(4.1), Inches(11), Inches(1.5))
    tf2 = txb2.text_frame
    tf2.word_wrap = True
    set_text(tf2, "Careers RAG Chatbot — Powered by Azure OpenAI + Elasticsearch", font_size=16, color=RGBColor(0xC0, 0xD0, 0xE0), alignment=PP_ALIGN.CENTER)
    add_paragraph(tf2, "MouriTech  •  April 2026", font_size=13, color=SLATE_500, space_before=Pt(20))
    tf2.paragraphs[1].alignment = PP_ALIGN.CENTER

    return slide


# ── End Slide ─────────────────────────────────────────────────────────────────

def create_end_slide(total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide)

    txb = add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1))
    tf = txb.text_frame
    set_text(tf, "Thank You", font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_paragraph(tf, "Questions & Discussion", font_size=22, color=RGBColor(0xA0, 0xC0, 0xE0), space_before=Pt(16))
    tf.paragraphs[1].alignment = PP_ALIGN.CENTER

    # Contact
    txb2 = add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(1))
    tf2 = txb2.text_frame
    set_text(tf2, "UNIDO AI Chatbot  •  MouriTech  •  April 2026", font_size=13, color=SLATE_500, alignment=PP_ALIGN.CENTER)

    return slide


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE CONTENT
# ══════════════════════════════════════════════════════════════════════════════

TOTAL_SLIDES = 18  # title + 16 content + end

SLIDES = [
    # ── 1. Architecture ──
    {
        "title": "Architecture",
        "subtitle": "Component overview and technology stack",
        "accent": UNIDO_BLUE,
        "bullets": [
            "Frontend — React 19 + Vite 7, Bootstrap 5, Recharts charts, Socket.IO Client, React Hook Form",
            "Backend — Express 5 on Node.js, Mongoose 9 (MongoDB ODM), Puppeteer (scraping), Socket.IO, node-cron",
            "AI Layer — Azure OpenAI GPT-4o-mini (chat completion) + text-embedding-3-small (1536-dim embeddings)",
            "Vector Store — Elasticsearch 8.x with kNN cosine similarity search on dense_vector index",
            "Database — MongoDB with 4 collections: chatSessions, chatLogs, adminUsers, adminSettings",
            "Security Stack — Helmet, strict CORS, 3-tier rate limits, JWT HS256, bcrypt 12 rounds, NoSQL injection guard",
            "Data Flow — User → React Widget → Express API → Guard Service → RAG Pipeline → JSON Response",
        ],
    },
    # ── 2. Chat Flow ──
    {
        "title": "Chat Flow",
        "subtitle": "End-to-end journey from user input to AI response",
        "accent": UNIDO_BLUE,
        "bullets": [
            "Session Init — GET /api/chat/visibility checks if chatbot is enabled; sessionId stored in sessionStorage (per-tab isolation)",
            "Message Path — POST /api/chat/ask → express-validator → enabled check → Guard checks → RAG → save to DB → return { answer, sources }",
            "Guard Gates — Two pre-RAG checks: isBlockedQuery() (offensive/confidential/injection) and isOutOfScope() (unrelated topics)",
            "History-Aware — Loads prior messages from ChatSession to give AI full conversation context",
            "Persist & Emit — Saves ChatLog with timing/IP/user-agent metadata; emits analytics:updated via Socket.IO for live admin dashboard",
            "Session Recovery — If backend returns 404, frontend auto-clears sessionStorage and creates a fresh session (zero user friction)",
            "Quick Actions — 3 preset buttons (\"What jobs are available?\", \"How do I apply?\", \"What are the benefits?\") use the same sendMessage() path",
        ],
    },
    # ── 3. RAG Pipeline ──
    {
        "title": "RAG Pipeline",
        "subtitle": "Retrieval-Augmented Generation: Embed → Retrieve → Generate",
        "accent": ACCENT_GREEN,
        "bullets": [
            "Step 1 — Embed: User question → Azure OpenAI text-embedding-3-small → 1536-dimensional float vector capturing semantic meaning",
            "Step 2 — Retrieve: Vector sent to Elasticsearch kNN search; TOP_K=5 results, num_candidates=50, minimum cosine similarity 0.75",
            "No Results Safeguard — If no chunks pass the 0.75 threshold, returns \"I do not have enough information\" — skips GPT call (saves cost, prevents hallucination)",
            "Step 3 — Generate: Retrieved chunks formatted as CONTEXT, sent with conversation history to GPT-4o-mini at temperature 0.1 (near-deterministic)",
            "System Prompt — Loaded from AdminSettings (configurable via Admin UI); instructs AI to answer ONLY from provided context",
            "Tunable via Env — RAG_TOP_K and RAG_MIN_SCORE configurable without code changes",
        ],
    },
    # ── 4. Guard Service ──
    {
        "title": "Guard Service",
        "subtitle": "Content filtering, prompt injection defense, and scope control",
        "accent": ACCENT_RED,
        "bullets": [
            "Four-Layer Filter — Every question passes: Offensive → Confidential → Prompt Injection → Out-of-Scope checks before reaching RAG",
            "Offensive — Blocks profanity, hate speech, threats, violence, self-harm keywords; returns professional fallback message",
            "Confidential — Blocks PII extraction: \"staff salary\", \"passport\", \"ssn\", \"credit card\", \"medical record\", \"employee data\"",
            "Prompt Injection — Blocks \"ignore previous instructions\", \"you are now\", \"system:\", \"DAN jailbreak\", \"[INST]\", \"<<SYS>>\", \"forget everything\" + 10 more",
            "Out-of-Scope — Polite redirection for code generation, recipes, weather, sports, stocks, crypto, homework requests",
            "Hard Limits — Empty inputs and messages over 2,000 characters blocked; all blocked queries logged with guardReason for audit",
            "Contextual Responses — buildBlockedAnswer(reason) returns different user-friendly messages per block type",
        ],
    },
    # ── 5. Data Pipeline ──
    {
        "title": "Data Pipeline",
        "subtitle": "Scraping, cleaning, chunking, embedding, and indexing",
        "accent": ACCENT_PURPLE,
        "bullets": [
            "6 Scrapers (Puppeteer) — Homepage, FAQs, Job Search (paginated 25/page), Job Details, Categories, Category Content",
            "HTML Cleaning — Cheerio strips <script>, <style>, <header>, <footer>, <nav> elements; 1.5s delay between pages",
            "Chunking — Sentence-based splitting, max 1000 chars/chunk, 1-sentence overlap for context continuity; 6 typed chunk formats",
            "Embedding — Each chunk → Azure OpenAI 1536-dim vector; processed sequentially to respect API rate limits",
            "Indexing — Elasticsearch dense_vector index with cosine similarity; deleteByQuery clears old data, then bulk API indexes all chunks",
            "Backup — JSON files saved to data/ and chunks/ with date-prefixed filenames (e.g., 180226_unido_faqs.json)",
            "Cron — Daily at midnight UTC (0 0 * * *) when ENABLE_CRON=true; also triggerable on-demand from Admin panel",
        ],
    },
    # ── 6. Auth Flow ──
    {
        "title": "Auth Flow",
        "subtitle": "Admin login, JWT verification, and route protection",
        "accent": ACCENT_GREEN,
        "bullets": [
            "Login — POST /api/admin/auth/login (rate limited 5/min) → bcrypt.compare (12 rounds) → JWT HS256 signed, 25-min expiry → stored in localStorage",
            "Timing Attack Defense — If user not found, a dummy password is hashed anyway to prevent user enumeration via response timing",
            "JWT Payload — { id, role: \"admin\", iat } with 32+ character secret required in production",
            "Middleware — Extracts Bearer token → validates 3-part JWT format → jwt.verify() with algorithms: [\"HS256\"] + maxAge: \"30m\" → checks role === \"admin\"",
            "Client-Side Guard — PrivateRoute.jsx parses JWT, checks exp with 30s buffer, validates role (UX only — server independently verifies every call)",
            "Error Handling — Expired tokens → specific 401; other failures → generic 401; role mismatch → 403 with IP logged",
        ],
    },
    # ── 7. Admin Ops ──
    {
        "title": "Admin Operations",
        "subtitle": "Analytics dashboard, settings, scraping control, exports, and data management",
        "accent": UNIDO_BLUE,
        "bullets": [
            "Analytics — MongoDB aggregations: conversations, message distribution, unique users (by IP), avg response time, error count; 3 Recharts charts",
            "Settings — PUT /api/admin/settings updates system prompt (max 500 chars) and chatbotEnabled toggle; propagates live via Socket.IO",
            "Scraping Control — POST /api/admin/scrape/trigger with state machine (idle → running); frontend polls status every 5s",
            "CSV Export — GET /api/admin/reports/chat-logs with date range + type filter; max 10,000 rows with timing/IP/status columns",
            "All Information — Paginated (default 25, max 100), regex-searchable across 6 fields, ReDoS-protected, sorted newest-first, auto-refreshes via Socket.IO",
        ],
    },
    # ── 8. Real-time ──
    {
        "title": "Real-time Events",
        "subtitle": "Socket.IO architecture for live updates across admin and widget",
        "accent": ACCENT_PURPLE,
        "bullets": [
            "Architecture — Socket.IO on same Express server; WebSocket transport with polling fallback; admin clients send JWT in auth handshake",
            "Auth Middleware — Verifies JWT on connect → sets socket.isAdmin = true; invalid/missing token still connects as non-admin",
            "Connection — Singleton pattern via getSocket(); server emits socket:ready with timestamp; 30-minute idle disconnect timeout",
            "analytics:updated — Emitted on session create, chat log save, or settings change → triggers admin dashboard refresh automatically",
            "information:updated — Emitted on chat log creation → triggers All Information table refresh in real-time",
            "chatbot:visibilityChanged — Emitted when admin toggles chatbot on/off → ChatWidget shows/hides instantly with { chatbotEnabled } payload",
        ],
    },
    # ── 9. Security ──
    {
        "title": "Security",
        "subtitle": "Defense-in-depth: 9 middleware layers protecting every request",
        "accent": ACCENT_RED,
        "bullets": [
            "Middleware Stack — requestId → securityHeaders → methodGuard → payloadGuard (2 MB) → Helmet + CORS → auditLogger → Body Parsers → inputSanitizer → noSQLInjectionGuard",
            "Rate Limits — Auth: 5/min (brute-force), Chat: 100/min, Admin: 30/min; Speed limiter: 500ms delay after 50 req/min",
            "Input Sanitizer — Blocks <script>, javascript:, on*=, ${}, {{}} patterns → 400 Bad Request",
            "NoSQL Guard — Rejects body/query keys starting with $ (MongoDB operators like $gt, $ne) → 400 Bad Request",
            "CSP — default-src 'self', script-src 'self', object-src 'none', frame-src 'none'; X-Powered-By removed",
            "CORS — Production: only https://careers.unido.org; blocks requests with no Origin header; strict allowed methods",
            "Headers — HSTS 1yr + preload, X-Frame-Options: DENY, nosniff, Cross-Origin policies: same-origin, HTTPS redirect in production",
        ],
    },
    # ── 10. API Reference ──
    {
        "title": "API Reference",
        "subtitle": "Complete list of REST endpoints (7 public + 8 protected)",
        "accent": UNIDO_BLUE,
        "bullets": [
            "GET /health — Service health check (public)",
            "GET /api/chat/visibility — Check if chatbot is enabled (public)",
            "POST /api/chat/session — Create anonymous chat session (public)",
            "POST /api/chat/ask — Send question, receive AI answer + sources (public, rate limited 100/min)",
            "POST /api/admin/auth/login — Admin login, returns JWT (public, rate limited 5/min)",
            "GET /api/admin/analytics — Dashboard metrics with optional date range filter (JWT required)",
            "GET /api/admin/allInformation — Paginated chat logs with search (JWT required)",
            "GET|PUT /api/admin/settings — Read/update system prompt and chatbot toggle (JWT required)",
            "POST /api/admin/scrape/trigger | GET /api/admin/scrape/status — Trigger and monitor scraping pipeline (JWT required)",
            "GET /api/admin/reports/chat-logs — CSV export with date/type filters (JWT required)",
        ],
    },
    # ── 11. Env Vars ──
    {
        "title": "Environment Variables",
        "subtitle": "All configuration values needed to run the system",
        "accent": SLATE_500,
        "bullets": [
            "Required — MONGO_URI, AZURE_OPENAI_ENDPOINT, AZURE_OPENAI_KEY, AZURE_OPENAI_DEPLOYMENT, AZURE_EMBEDDING_DEPLOYMENT, JWT_SECRET, ADMIN_EMAIL, ADMIN_PASSWORD",
            "Elasticsearch — ELASTIC_NODE (localhost:9200), ELASTIC_USERNAME, ELASTIC_PASSWORD",
            "Azure OpenAI — AZURE_OPENAI_API_VERSION (default: 2024-10-21)",
            "Tunable — JWT_EXPIRES (25m), ENABLE_CRON (false), CRON_TIMEZONE (UTC), RAG_TOP_K (5), RAG_MIN_SCORE (0.75)",
            "Frontend — Single variable: VITE_API_BASE_URL (default: http://localhost:5000)",
            "Production Enforcement — NODE_ENV=production requires JWT_SECRET 32+ chars; ADMIN_PASSWORD 12+ chars recommended",
            "Security — Never commit .env to source control; use Azure Key Vault or secrets manager in production",
        ],
    },
    # ── 12. Org Security (Part 1) ──
    {
        "title": "Org Security — Data & Auth",
        "subtitle": "Data protection, privacy, authentication hardening, and network security",
        "accent": ACCENT_RED,
        "bullets": [
            "No PII Storage — Anonymous UUID sessions; no personal data collected; GDPR / UN data-protection compliant by design",
            "PII Query Blocking — Guard blocks attempts to extract \"staff salary\", \"passport\", \"ssn\", \"credit card\" via the AI",
            "Response Grounding — RAG-only answers at temperature 0.1; no chunks = no GPT call (prevents hallucination entirely)",
            "Session Isolation — Per-tab sessionStorage; sessions can't cross-reference; closing tab destroys session data",
            "Auth Hardening — bcrypt 12 rounds (~250ms/hash), timing-attack defense (dummy hash on user not found), JWT HS256 with 25m expiry",
            "Auth Rate Limiting — 5 login attempts/min per IP; does NOT skip successful requests — prevents credential stuffing",
            "Network Security — HSTS 1yr + preload, strict CORS, X-Frame-Options: DENY, Permissions-Policy disables camera/mic/geo",
        ],
    },
    # ── 13. Org Security (Part 2) ──
    {
        "title": "Org Security — Input & OWASP",
        "subtitle": "Input security, AI safety, monitoring, and OWASP Top 10 compliance",
        "accent": ACCENT_RED,
        "bullets": [
            "XSS Defense — 3-layer: Input sanitizer middleware + Content Security Policy + React's default output escaping",
            "NoSQL Injection — Dedicated guard rejects $-prefixed keys ($gt, $ne, $or); template injection (${}, {{}}) also blocked",
            "Prompt Injection — Four-layer guard filters before AI: blocks system prompt overrides, DAN jailbreaks, instruction manipulation",
            "AI-Specific Risks — Hallucination: RAG + low temp; Data poisoning: only official UNIDO domains scraped; Scope creep: out-of-scope detector",
            "Monitoring — Unique requestId per request; audit logger records method/path/status/duration/IP; errors ≥400 and slow >5s requests flagged",
            "OWASP Top 10 — 9/10 categories \"Mitigated\" (A01–A05, A07–A10); A06 \"Ongoing\" — modern deps, regular npm audit recommended",
            "No SSRF — No user-supplied URLs fetched; scraping targets hardcoded to UNIDO domains only",
        ],
    },
    # ── 14. Go-Live (Part 1) ──
    {
        "title": "Go-Live — Azure OpenAI & Elasticsearch",
        "subtitle": "The AI intelligence and semantic search backbone",
        "accent": UNIDO_BLUE,
        "bullets": [
            "Azure OpenAI API [CRITICAL] — Powers both question understanding (embeddings) and human-like response generation (GPT-4o-mini)",
            "Credentials Needed — Endpoint URL, API key, chat model deployment name, embedding model deployment name, API version",
            "Env Vars — AZURE_OPENAI_ENDPOINT, AZURE_OPENAI_KEY, AZURE_OPENAI_DEPLOYMENT, AZURE_EMBEDDING_DEPLOYMENT, AZURE_OPENAI_API_VERSION",
            "Elasticsearch [CRITICAL] — Vector database enabling fast kNN semantic search across all scraped career content (1536-dim cosine)",
            "ES Credentials — Node URL, username, password, index name; requires Elasticsearch 8.x+ with dense_vector + kNN support",
            "ES Options — Elastic Cloud, Azure-hosted Elasticsearch, or self-managed cluster",
        ],
    },
    # ── 15. Go-Live (Part 2) ──
    {
        "title": "Go-Live — Blob Storage & Content Safety",
        "subtitle": "Persistent storage, backups, and ML-powered guardrails",
        "accent": ACCENT_AMBER,
        "bullets": [
            "Azure Blob Storage [CRITICAL] — Stores raw scraped JSON, chunked content, and backups between scraping and indexing pipelines",
            "Re-indexing — If Elasticsearch needs rebuilding (schema change, corruption), system re-reads from Blob Storage without re-scraping the website",
            "Blob Credentials — AZURE_STORAGE_CONNECTION_STRING + container name (default: unido-chatbot)",
            "Azure AI Content Safety [IMPORTANT] — ML-powered guardrails scanning responses for harmful content; provides auditable safety scores",
            "Complements Built-in Guard — Local regex guard catches known patterns instantly; Azure Content Safety catches novel threats with continuously-updated ML models",
            "Content Safety Credentials — AZURE_CONTENT_SAFETY_ENDPOINT + AZURE_CONTENT_SAFETY_KEY",
        ],
    },
    # ── 16. Go-Live (Part 3) ──
    {
        "title": "Go-Live — MongoDB & Summary",
        "subtitle": "Operational database and complete service checklist",
        "accent": ACCENT_GREEN,
        "bullets": [
            "MongoDB [CRITICAL] — Primary operational database storing chat sessions, chat logs (full audit trail), admin users, admin settings, scraped content metadata",
            "Credential — Single connection string (MONGO_URI); requires MongoDB 6+; options: Azure Cosmos DB, MongoDB Atlas, or self-managed",
            "Data Retention — Establish a purge policy (e.g., 90 days) for chat logs; currently no automatic cleanup configured",
            "Service #1: Azure OpenAI API — Intelligence layer (embeddings + chat) [CRITICAL]",
            "Service #2: Elasticsearch — Semantic vector search backbone [CRITICAL]",
            "Service #3: Azure Blob Storage — Raw data storage and backup [CRITICAL]",
            "Service #4: Azure AI Content Safety — ML-powered content guardrails [IMPORTANT]",
            "Service #5: MongoDB — Operational data, logs, and analytics [CRITICAL]",
        ],
    },
]


# ══════════════════════════════════════════════════════════════════════════════
#  BUILD PRESENTATION
# ══════════════════════════════════════════════════════════════════════════════

print("Generating UNIDO AI Chatbot slide deck...")

# Title slide
create_title_slide()

# Content slides
for i, slide_data in enumerate(SLIDES, start=1):
    create_content_slide(
        slide_num=i,
        total=len(SLIDES),
        title=slide_data["title"],
        subtitle=slide_data["subtitle"],
        bullets=slide_data["bullets"],
        accent_color=slide_data.get("accent", UNIDO_BLUE),
    )

# End slide
create_end_slide(len(SLIDES))

# Save
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "docs", "UNIDO_AI_Chatbot_Presentation.pptx")
os.makedirs(os.path.dirname(out_path), exist_ok=True)
prs.save(out_path)
print(f"Saved to: {out_path}")
print(f"Total slides: {len(prs.slides)} (1 title + {len(SLIDES)} content + 1 closing)")
